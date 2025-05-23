# app.py

"""
==============
A minimal RAG back-end for your Streamlit client.

▶︎ Quick-start
pip install fastapi uvicorn boto3 sentence_transformers faiss-cpu PyMuPDF
           python-multipart tavily-python pillow pandas python-docx python-pptx

▶︎ Run
uvicorn app:app --host 0.0.0.0 --port 8000
"""

import os, io, json, uuid, tempfile, pickle, fitz, faiss, boto3
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from datetime import datetime
import numpy as np
import pandas as pd
from fastapi import FastAPI, UploadFile, File, HTTPException
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
from docx import Document
from pptx import Presentation
from PIL import Image
from tavily import TavilyClient              # pip install tavily-python

# ── CONFIG ──────────────────────────────────────────────────────────────────
REGION   = "us-east-1"
BUCKET   = "satagroup-test"
TAVILY   = "YOUR-TAVILY-KEY"
MODEL_ID = "anthropic.claude-3-7-sonnet-20250219-v1:0"
INDEX_F  = "faiss_index.bin"
META_F   = "metadata_store.pkl"
EMB_DIM  = 768

# ── AWS clients (re-use across calls) ───────────────────────────────────────
s3       = boto3.client("s3",      region_name=REGION)
textract = boto3.client("textract",region_name=REGION)
bedrock  = boto3.client("bedrock-runtime", region_name=REGION)

# ── Globals (in-memory cache) ───────────────────────────────────────────────
app          = FastAPI(title="Document-RAG Gateway", version="1.0.0")
mpnet_model  = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")
faiss_index  : faiss.Index  # Will be loaded at startup
metadata     : List[Dict]   # Will be loaded at startup

# ── Helpers ─────────────────────────────────────────────────────────────────
def _load_stores() -> None:
    """Load FAISS + metadata from S3 (or initialise fresh)."""
    global faiss_index, metadata
    try:
        s3.download_file(BUCKET, INDEX_F, INDEX_F)
        s3.download_file(BUCKET, META_F,  META_F)
        faiss_index = faiss.read_index(INDEX_F)
        metadata    = pickle.load(open(META_F, "rb"))
    except Exception:
        faiss_index = faiss.IndexFlatL2(EMB_DIM)         # empty
        metadata    = []

def _persist_stores() -> None:
    faiss.write_index(faiss_index, INDEX_F)
    pickle.dump(metadata, open(META_F, "wb"))
    s3.upload_file(INDEX_F, BUCKET, INDEX_F)
    s3.upload_file(META_F, BUCKET, META_F)

def _embed(txt: str) -> np.ndarray:
    return mpnet_model.encode(txt, normalize_embeddings=True)

def _ocr_bytes(blob: bytes) -> str:
    """Use Amazon Textract's detect_document_text for basic OCR."""
    resp = textract.detect_document_text(Document={'Bytes': blob})
    return "\n".join(b["Text"] for b in resp.get("Blocks", [])
                     if b["BlockType"] == "LINE")

def _process_file(path: Path, owner: str) -> None:
    """Extract text, embed, push into FAISS & metadata."""
    ext = path.suffix.lower()
    chunks: List[Tuple[str,int]] = []   # [(text,page)]

    if ext == ".pdf":
        doc = fitz.open(path)
        for p in doc:
            # Convert page to PNG and run OCR
            png_bytes = p.get_pixmap().tobytes("png")
            text = _ocr_bytes(png_bytes)
            chunks.append((text, p.number + 1))

    elif ext in {".jpg", ".jpeg", ".png"}:
        # Single image
        img_bytes = path.read_bytes()
        text = _ocr_bytes(img_bytes)
        chunks.append((text, 1))

    elif ext in {".doc", ".docx"}:
        full = "\n".join(p.text for p in Document(path).paragraphs if p.text)
        # Simple chunking for docx
        step = 1000
        for i in range(0, len(full), step):
            page_text = full[i:i+step]
            page_num = (i // step) + 1
            chunks.append((page_text, page_num))

    elif ext == ".pptx":
        prs = Presentation(path)
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = "\n".join(s.text for s in slide.shapes if hasattr(s, "text"))
            chunks.append((slide_text, idx))

    elif ext in {".csv", ".xlsx"}:
        if ext == ".csv":
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
        # chunk every 50 rows
        for i in range(0, len(df), 50):
            chunk_text = df.iloc[i:i+50].to_string(index=False)
            page_num = (i // 50) + 1
            chunks.append((chunk_text, page_num))

    else:
        raise ValueError(f"Unsupported file type: {ext}")

    # embed & store
    for txt, page in chunks:
        vec = _embed(txt)
        faiss_index.add(vec.reshape(1, -1))
        metadata.append({
            "filename": path.name,
            "page": page,
            "text": txt,
            "owner": owner,
            "uploaded": datetime.utcnow().isoformat()
        })

# ── Pydantic DTOs ───────────────────────────────────────────────────────────
class UploadAck(BaseModel):
    status: str
    filename: str
    pages_indexed: int

class QueryBody(BaseModel):
    selected_files: List[str]
    selected_page_ranges: Dict[str, Tuple[int,int]]
    prompt: str
    top_k: int = 20
    last_messages: List[str] = []
    web_search: bool = False

class Answer(BaseModel):
    answer: str
    context: List[Dict]

# ── Start-up ────────────────────────────────────────────────────────────────
@app.on_event("startup")
def _warm():
    _load_stores()

# ── End-points ───────────────────────────────────────────────────────────────
@app.post("/upload_document", response_model=UploadAck)
async def upload_document(file: UploadFile = File(...),
                          owner: str = "anonymous"):
    """
    Upload a single document (PDF, image, DOCX, PPTX, XLSX, CSV) and index it synchronously.
    """
    fn = file.filename
    if not fn:
        raise HTTPException(400, "No filename")

    # Save to a temp file
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=Path(fn).suffix)
    tmp.write(await file.read())
    tmp.close()

    # 1) store original in S3 so users can download
    with open(tmp.name, "rb") as fh:
        s3.upload_fileobj(fh, BUCKET, fn)

    # 2) process & embed synchronously
    try:
        _process_file(Path(tmp.name), owner)  # extracts text, embed, store in FAISS
        _persist_stores()                    # writes & uploads index/metadata
    finally:
        os.remove(tmp.name)

    # 3) count how many pages/chunks we just added
    pages = sum(1 for m in metadata if m["filename"] == fn)
    return UploadAck(status="done", filename=fn, pages_indexed=pages)

@app.post("/query_documents_with_page_range", response_model=Answer)
def query_docs(body: QueryBody):
    """
    Given a user prompt, a list of selected files + page ranges, returns top-K hits + LLM answer.
    """
    if faiss_index.ntotal == 0:
        raise HTTPException(404, "Index empty")

    # embed the user's prompt
    q_vec = _embed(body.prompt).reshape(1, -1)
    # oversample, we'll filter
    D, I = faiss_index.search(q_vec, min(body.top_k * 5, faiss_index.ntotal))

    hits = []
    for dist, idx in zip(D[0], I[0]):
        meta = metadata[idx]
        fname = meta["filename"]
        pg = meta["page"]
        # filter by user selected
        if body.selected_files and fname not in body.selected_files:
            continue
        # filter by page range
        lo, hi = body.selected_page_ranges.get(fname, (1, 10**6))
        if not (lo <= pg <= hi):
            continue

        hits.append({"score": float(dist), **meta})
        if len(hits) >= body.top_k:
            break

    tav_results = {}
    if body.web_search and TAVILY:
        cli = TavilyClient(api_key=TAVILY)
        tav_results = cli.search(body.prompt, search_depth="advanced",
                                 include_raw_content=True)

    # build final LLM prompt
    ctx = {"hits": hits, "last": body.last_messages, "tavily": tav_results}
    doc_context = json.dumps(ctx, indent=2)

    payload = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": body.prompt + "\n\n# Context:\n" + doc_context
                    }
                ]
            }
        ]
    }

    out = bedrock.invoke_model(
        modelId=MODEL_ID,
        contentType="application/json",
        accept="application/json",
        body=json.dumps(payload)
    )
    # read the response
    answer = json.loads(out["body"].read())["content"][0]["text"]
    return Answer(answer=answer, context=hits)
