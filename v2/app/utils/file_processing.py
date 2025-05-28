"""
File processing utilities for text extraction from different file types.
"""
import os
import fitz  # PyMuPDF
import pandas as pd
from pathlib import Path
from typing import List, Tuple, Dict, Any
from docx import Document
from pptx import Presentation

from app.core.aws_client import aws_client

def extract_text_from_pdf(file_path: str) -> List[Tuple[str, int]]:
    """
    Extract text from PDF file.
    
    Args:
        file_path: Path to PDF file
        
    Returns:
        List of (text, page_number) tuples
    """
    chunks = []
    try:
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc):
            # Get high-resolution image for better OCR
            png = page.get_pixmap(dpi=300).tobytes("png")
            text = aws_client.extract_text_from_image(png)
            if text.strip():
                chunks.append((text, page_num + 1))
        return chunks
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return []

def extract_text_from_image(file_path: str) -> List[Tuple[str, int]]:
    """
    Extract text from image file.
    
    Args:
        file_path: Path to image file
        
    Returns:
        List of (text, page_number) tuples
    """
    chunks = []
    try:
        with open(file_path, "rb") as f:
            image_bytes = f.read()
        text = aws_client.extract_text_from_image(image_bytes)
        if text.strip():
            chunks.append((text, 1))
        return chunks
    except Exception as e:
        print(f"Error extracting text from image: {e}")
        return []

def extract_text_from_docx(file_path: str) -> List[Tuple[str, int]]:
    """
    Extract text from DOCX file.
    
    Args:
        file_path: Path to DOCX file
        
    Returns:
        List of (text, chunk_number) tuples
    """
    chunks = []
    try:
        doc = Document(file_path)
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text)
        
        # Split into chunks of 1000 characters
        for i in range(0, len(full_text), 1000):
            chunk = full_text[i:i+1000]
            if chunk.strip():
                chunks.append((chunk, i//1000 + 1))
        return chunks
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return []

def extract_text_from_pptx(file_path: str) -> List[Tuple[str, int]]:
    """
    Extract text from PPTX file.
    
    Args:
        file_path: Path to PPTX file
        
    Returns:
        List of (text, slide_number) tuples
    """
    chunks = []
    try:
        prs = Presentation(file_path)
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = "\n".join(
                s.text for s in slide.shapes if hasattr(s, "text")
            )
            if slide_text.strip():
                chunks.append((slide_text, idx))
        return chunks
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return []

def extract_text_from_spreadsheet(file_path: str) -> List[Tuple[str, int]]:
    """
    Extract text from CSV or Excel file.
    
    Args:
        file_path: Path to CSV or Excel file
        
    Returns:
        List of (text, chunk_number) tuples
    """
    chunks = []
    try:
        ext = os.path.splitext(file_path)[1].lower()
        df = pd.read_csv(file_path) if ext == ".csv" else pd.read_excel(file_path)
        
        # Split into chunks of 50 rows
        for i in range(0, len(df), 50):
            chunk = df.iloc[i:i+50].to_string(index=False)
            if chunk.strip():
                chunks.append((chunk, i//50 + 1))
        return chunks
    except Exception as e:
        print(f"Error extracting text from spreadsheet: {e}")
        return []

def process_file(file_path: str, owner: str, original_filename: str) -> int:
    """
    Process file and extract text chunks.
    
    Args:
        file_path: Path to file
        owner: Owner of the file
        original_filename: Original filename
        
    Returns:
        Number of chunks processed
    """
    from app.core.embedding import embedding_manager
    
    ext = os.path.splitext(file_path)[1].lower()
    chunks = []
    
    # Extract text based on file type
    if ext == ".pdf":
        chunks = extract_text_from_pdf(file_path)
    elif ext in {".jpg", ".jpeg", ".png"}:
        chunks = extract_text_from_image(file_path)
    elif ext in {".doc", ".docx"}:
        chunks = extract_text_from_docx(file_path)
    elif ext == ".pptx":
        chunks = extract_text_from_pptx(file_path)
    elif ext in {".csv", ".xlsx"}:
        chunks = extract_text_from_spreadsheet(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    
    # Add chunks to embedding index
    for text, page in chunks:
        embedding_manager.add_to_index(text, {
            "filename": original_filename,
            "page": page,
            "owner": owner
        })
    
    return len(chunks)
