"""
FastAPI Skill Analysis System
Based on existing Streamlit architecture with FAISS indexing and AWS integration
"""

import os
import json
import tempfile
import pickle
import fitz  # PyMuPDF
import faiss
import boto3
import uuid
import hashlib
import io
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from datetime import datetime, timedelta
import numpy as np
import pandas as pd
from fastapi import FastAPI, UploadFile, File, HTTPException, Depends, status
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel, EmailStr
from sentence_transformers import SentenceTransformer
from docx import Document
from pptx import Presentation
import tiktoken

# ── CONFIG ──────────────────────────────────────────────────────────────────
REGION = "us-east-1"
BUCKET = "satagroup-test"
MODEL_ID = "arn:aws:bedrock:us-east-1:343218220592:inference-profile/us.anthropic.claude-3-7-sonnet-20250219-v1:0"
INDEX_F = "faiss_index.bin"
META_F = "metadata_store.pkl"
USERS_F = "users.json"
ANALYSIS_F = "analysis_results.json"
CHAT_F = "chat_conversations.json"
EMB_DIM = 768

# AWS Credentials


# ── AWS clients ─────────────────────────────────────────────────────────────
s3 = boto3.client(
    "s3", 
    region_name=REGION
)
textract = boto3.client(
    "textract", 
    region_name=REGION
)
bedrock = boto3.client(
    "bedrock-runtime", 
    region_name=REGION
)

# ── Globals ─────────────────────────────────────────────────────────────────
app = FastAPI(title="Skill Analysis API", version="1.0.0", description="HR-tech skill analysis and document management system")
embed_model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")
security = HTTPBearer()

# Global variables for FAISS and metadata
faiss_index: faiss.Index
metadata: List[Dict]
users_data: Dict
analysis_results: Dict
chat_conversations: Dict

# ── Data Models ─────────────────────────────────────────────────────────────
class UserRegister(BaseModel):
    fullName: str
    email: EmailStr
    password: str
    role: str = "user"  # user|manager|admin

class UserLogin(BaseModel):
    email: EmailStr
    password: str

class ForgotPassword(BaseModel):
    email: EmailStr

class ResetPassword(BaseModel):
    token: str
    newPassword: str

class AnalysisRequest(BaseModel):
    jdId: str
    cvIds: List[str]
    options: Dict = {"includeScores": True, "language": "en"}

class UploadResponse(BaseModel):
    jdId: Optional[str] = None
    cvId: Optional[str] = None
    title: Optional[str] = None
    fileName: str
    uploadedAt: str

class AnalysisResponse(BaseModel):
    analysisId: str
    timestamp: str
    results: Dict  # Changed from str to Dict to support JSON format

class ChatRequest(BaseModel):
    jdId: str
    cvId: str
    message: str
    conversationId: Optional[str] = None

class ChatResponse(BaseModel):
    conversationId: str
    response: str
    timestamp: str

class ChatMessage(BaseModel):
    role: str  # "user" or "assistant"
    content: str
    timestamp: str

class ChatConversation(BaseModel):
    conversationId: str
    jdId: str
    cvId: str
    userId: str
    messages: List[ChatMessage]
    createdAt: str
    updatedAt: str

# ── Utility Functions ───────────────────────────────────────────────────────
def init_files():
    """Initialize required files if they don't exist"""
    # Initialize users.json
    if not os.path.exists(USERS_F):
        default_users = {
            "admin": {
                "fullName": "System Admin",
                "email": "admin@system.com",
                "password": "admin123",
                "role": "admin",
                "userId": str(uuid.uuid4()),
                "createdAt": datetime.utcnow().isoformat()
            }
        }
        with open(USERS_F, "w") as f:
            json.dump(default_users, f, indent=2)
    
    # Initialize analysis results
    if not os.path.exists(ANALYSIS_F):
        with open(ANALYSIS_F, "w") as f:
            json.dump({}, f, indent=2)
    
    # Initialize chat conversations
    if not os.path.exists(CHAT_F):
        with open(CHAT_F, "w") as f:
            json.dump({}, f, indent=2)

def load_users():
    """Load users from JSON file"""
    global users_data
    try:
        with open(USERS_F, "r") as f:
            users_data = json.load(f)
    except FileNotFoundError:
        users_data = {}

def save_users():
    """Save users to JSON file"""
    with open(USERS_F, "w") as f:
        json.dump(users_data, f, indent=2)
    # Also upload to S3
    try:
        s3.upload_file(USERS_F, BUCKET, USERS_F)
    except Exception as e:
        print(f"Warning: Could not upload users.json to S3: {e}")

def load_analysis_results():
    """Load analysis results from JSON file"""
    global analysis_results
    try:
        with open(ANALYSIS_F, "r") as f:
            analysis_results = json.load(f)
    except FileNotFoundError:
        analysis_results = {}

def save_analysis_results():
    """Save analysis results to JSON file"""
    with open(ANALYSIS_F, "w") as f:
        json.dump(analysis_results, f, indent=2)
    # Also upload to S3
    try:
        s3.upload_file(ANALYSIS_F, BUCKET, ANALYSIS_F)
    except Exception as e:
        print(f"Warning: Could not upload analysis_results.json to S3: {e}")

def load_chat_conversations():
    """Load chat conversations from JSON file"""
    global chat_conversations
    try:
        with open(CHAT_F, "r") as f:
            chat_conversations = json.load(f)
    except FileNotFoundError:
        chat_conversations = {}

def save_chat_conversations():
    """Save chat conversations to JSON file and S3"""
    with open(CHAT_F, "w") as f:
        json.dump(chat_conversations, f, indent=2)
    # Also upload to S3
    try:
        s3.upload_file(CHAT_F, BUCKET, CHAT_F)
    except Exception as e:
        print(f"Warning: Could not upload chat_conversations.json to S3: {e}")

def generate_token(user_data: dict) -> str:
    """Generate a simple token for user authentication"""
    payload = f"{user_data['email']}:{user_data['userId']}"
    return hashlib.sha256(payload.encode()).hexdigest()

def verify_token(token: str) -> Optional[dict]:
    """Verify token and return user data"""
    # Simple token verification - in production, use proper JWT
    for username, user_data in users_data.items():
        expected_token = generate_token(user_data)
        if token == expected_token:
            return user_data
    return None

def get_current_user(credentials: HTTPAuthorizationCredentials = Depends(security)) -> dict:
    """Get current authenticated user"""
    token = credentials.credentials
    user = verify_token(token)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Invalid authentication credentials",
            headers={"WWW-Authenticate": "Bearer"},
        )
    return user

def _load_stores() -> None:
    """Load FAISS index and metadata from storage"""
    global faiss_index, metadata
    try:
        # Try to download from S3 first
        s3.download_file(BUCKET, INDEX_F, INDEX_F)
        s3.download_file(BUCKET, META_F, META_F)
        faiss_index = faiss.read_index(INDEX_F)
        with open(META_F, "rb") as f:
            metadata = pickle.load(f)
        print("Index and metadata loaded from S3.")
    except Exception as e:
        print(f"Could not load from S3: {e}. Initializing new index.")
        faiss_index = faiss.IndexFlatL2(EMB_DIM)
        metadata = []

def _persist_stores() -> None:
    """Save FAISS index and metadata to storage"""
    faiss.write_index(faiss_index, INDEX_F)
    with open(META_F, "wb") as f:
        pickle.dump(metadata, f)
    # Upload to S3
    try:
        s3.upload_file(INDEX_F, BUCKET, INDEX_F)
        s3.upload_file(META_F, BUCKET, META_F)
    except Exception as e:
        print(f"Warning: Could not upload to S3: {e}")

def _embed(text: str) -> np.ndarray:
    """Generate embeddings for text"""
    return embed_model.encode(text, normalize_embeddings=True)

def _ocr_png(png_bytes: bytes) -> str:
    """Extract text from PNG using AWS Textract"""
    try:
        resp = textract.detect_document_text(Document={"Bytes": png_bytes})
        return "\n".join(b["Text"] for b in resp.get("Blocks", [])
                        if b["BlockType"] == "LINE")
    except Exception as e:
        print(f"OCR error: {e}")
        return ""

# ── Startup Event ──────────────────────────────────────────────────────────
@app.on_event("startup")
async def startup_event():
    """Initialize the application"""
    init_files()
    load_users()
    load_analysis_results()
    load_chat_conversations()
    _load_stores()
    print("Skill Analysis API started successfully!")

# ── Authentication Endpoints ───────────────────────────────────────────────
@app.post("/v1/auth/register")
async def register_user(user: UserRegister):
    """Register a new user"""
    # Check if user already exists
    if any(u.get("email") == user.email for u in users_data.values()):
        raise HTTPException(status_code=400, detail="Email already registered")
    
    # Create new user
    user_id = str(uuid.uuid4())
    username = user.email.split("@")[0]  # Use email prefix as username
    
    # Ensure unique username
    counter = 1
    original_username = username
    while username in users_data:
        username = f"{original_username}_{counter}"
        counter += 1
    
    users_data[username] = {
        "userId": user_id,
        "fullName": user.fullName,
        "email": user.email,
        "password": user.password,  # Plain text as requested
        "role": user.role,
        "createdAt": datetime.utcnow().isoformat()
    }
    
    save_users()
    
    return {
        "userId": user_id,
        "fullName": user.fullName,
        "email": user.email,
        "role": user.role
    }

@app.post("/v1/auth/login")
async def login_user(credentials: UserLogin):
    """Authenticate user and return token"""
    # Find user by email
    user_data = None
    for username, data in users_data.items():
        if data.get("email") == credentials.email and data.get("password") == credentials.password:
            user_data = data
            break
    
    if not user_data:
        raise HTTPException(status_code=401, detail="Invalid email or password")
    
    # Generate token
    access_token = generate_token(user_data)
    
    return {
        "accessToken": access_token,
        "expiresIn": 86400,  # 24 hours
        "refreshToken": access_token,  # Same for simplicity
        "user": {
            "userId": user_data["userId"],
            "fullName": user_data["fullName"],
            "email": user_data["email"],
            "role": user_data["role"]
        }
    }

@app.post("/v1/auth/forgot-password")
async def forgot_password(request: ForgotPassword):
    """Send password reset link (mock implementation)"""
    # Check if email exists
    email_exists = any(u.get("email") == request.email for u in users_data.values())
    
    # Always return success for security (don't reveal if email exists)
    return {
        "message": "If an account with this email exists, a password reset link has been sent."
    }

@app.post("/v1/auth/reset-password")
async def reset_password(request: ResetPassword):
    """Reset password using token (mock implementation)"""
    # In a real implementation, you would validate the reset token
    # For now, we'll just return a success message
    return {
        "message": "Password has been reset successfully."
    }

# ── Document Processing Functions ──────────────────────────────────────────
def _process_file(file_path: Path, owner: str, file_name: str, doc_type: str = "general") -> int:
    """Process uploaded file and add to FAISS index. Returns number of chunks processed."""
    ext = file_path.suffix.lower()
    chunks: List[Tuple[str, int]] = []
    
    if ext == ".pdf":
        doc = fitz.open(file_path)
        for pg in doc:
            # Higher DPI for better OCR
            png = pg.get_pixmap(dpi=300).tobytes("png")
            text = _ocr_png(png)
            if text.strip():
                chunks.append((text, pg.number + 1))
    
    elif ext in {".jpg", ".jpeg", ".png"}:
        text = _ocr_png(file_path.read_bytes())
        if text.strip():
            chunks.append((text, 1))
    
    elif ext in {".doc", ".docx"}:
        doc = Document(file_path)
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text)
        # Split into chunks of ~1000 characters
        for i in range(0, len(full_text), 1000):
            chunk = full_text[i:i+1000]
            if chunk.strip():
                chunks.append((chunk, i//1000 + 1))
    
    elif ext == ".pptx":
        prs = Presentation(file_path)
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
            if slide_text.strip():
                chunks.append((slide_text, idx))
    
    elif ext in {".csv", ".xlsx"}:
        df = pd.read_csv(file_path) if ext == ".csv" else pd.read_excel(file_path)
        for i in range(0, len(df), 50):
            chunk = df.iloc[i:i+50].to_string(index=False)
            if chunk.strip():
                chunks.append((chunk, i//50 + 1))
    else:
        raise HTTPException(400, f"Unsupported file type: {ext}")
    
    # Add chunks to FAISS index
    for text, page in chunks:
        embedding = _embed(text).reshape(1, -1)
        faiss_index.add(embedding)
        metadata.append({
            "filename": file_name,
            "page": page,
            "text": text,
            "owner": owner,
            "doc_type": doc_type,  # "jd" or "cv" or "general"
            "uploaded": datetime.utcnow().isoformat(),
            "doc_id": str(uuid.uuid4())
        })
    
    return len(chunks)

# ── Document Management Endpoints ──────────────────────────────────────────
@app.get("/v1/jds")
async def list_job_descriptions(current_user: dict = Depends(get_current_user)):
    """List all Job Descriptions"""
    jds = []
    seen_files = set()
    
    for item in metadata:
        if (item.get("doc_type") == "jd" and 
            item["filename"] not in seen_files and
            (item.get("owner") == current_user.get("email") or current_user.get("role") == "admin")):
            jds.append({
                "jdId": item.get("doc_id", str(uuid.uuid4())),
                "title": item["filename"].replace(".pdf", "").replace(".docx", ""),
                "uploadedAt": item.get("uploaded", datetime.utcnow().isoformat())
            })
            seen_files.add(item["filename"])
    
    return jds

@app.post("/v1/jds")
async def upload_job_description(
    file: UploadFile = File(...),
    title: Optional[str] = None,
    current_user: dict = Depends(get_current_user)
):
    """Upload a new Job Description"""
    if not file.filename:
        raise HTTPException(400, "No filename provided")
    
    # Validate file type
    allowed_extensions = {".pdf", ".docx", ".doc"}
    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in allowed_extensions:
        raise HTTPException(400, f"Unsupported file type. Allowed: {allowed_extensions}")
    
    # Create temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = Path(tmp.name)
    
    try:
        # Upload original file to S3
        s3.upload_fileobj(
            io.BytesIO(content),
            BUCKET,
            f"jds/{file.filename}"
        )
        
        # Process and index the file
        pages_indexed = _process_file(tmp_path, current_user["email"], file.filename, "jd")
        _persist_stores()
        
        jd_id = str(uuid.uuid4())
        display_title = title or file.filename.replace(".pdf", "").replace(".docx", "")
        
        return UploadResponse(
            jdId=jd_id,
            title=display_title,
            fileName=file.filename,
            uploadedAt=datetime.utcnow().isoformat()
        )
    
    finally:
        # Clean up temporary file
        os.unlink(tmp_path)

@app.get("/v1/jds/{jdId}")
async def get_job_description(jdId: str, current_user: dict = Depends(get_current_user)):
    """Get details of a specific Job Description"""
    # Find JD in metadata
    jd_content = []
    jd_info = None
    
    for item in metadata:
        if (item.get("doc_type") == "jd" and 
            (item.get("doc_id") == jdId or item["filename"].replace(".pdf", "").replace(".docx", "") == jdId)):
            if not jd_info:
                jd_info = item
            jd_content.append(item["text"])
    
    if not jd_info:
        raise HTTPException(404, "Job Description not found")
    
    # Check access permissions
    if (jd_info.get("owner") != current_user.get("email") and 
        current_user.get("role") != "admin"):
        raise HTTPException(403, "Access denied")
    
    return {
        "jdId": jdId,
        "title": jd_info["filename"].replace(".pdf", "").replace(".docx", ""),
        "content": "\n\n".join(jd_content),
        "uploadedAt": jd_info.get("uploaded", datetime.utcnow().isoformat())
    }

@app.get("/v1/cvs")
async def list_cvs(current_user: dict = Depends(get_current_user)):
    """List all uploaded CVs"""
    cvs = []
    seen_files = set()
    
    for item in metadata:
        if (item.get("doc_type") == "cv" and 
            item["filename"] not in seen_files and
            (item.get("owner") == current_user.get("email") or current_user.get("role") == "admin")):
            cvs.append({
                "cvId": item.get("doc_id", str(uuid.uuid4())),
                "fileName": item["filename"],
                "level": item.get("level", "unknown"),
                "uploadedAt": item.get("uploaded", datetime.utcnow().isoformat())
            })
            seen_files.add(item["filename"])
    
    return cvs

@app.post("/v1/cvs")
async def upload_cv(
    file: UploadFile = File(...),
    level: str = "jr",  # jr|mid|sr
    current_user: dict = Depends(get_current_user)
):
    """Upload a CV with experience level"""
    if not file.filename:
        raise HTTPException(400, "No filename provided")
    
    # Validate file type
    allowed_extensions = {".pdf", ".docx", ".doc"}
    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in allowed_extensions:
        raise HTTPException(400, f"Unsupported file type. Allowed: {allowed_extensions}")
    
    # Validate level
    if level not in ["jr", "mid", "sr"]:
        raise HTTPException(400, "Level must be one of: jr, mid, sr")
    
    # Create temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=file_ext) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = Path(tmp.name)
    
    try:
        # Upload original file to S3
        s3.upload_fileobj(
            io.BytesIO(content),
            BUCKET,
            f"cvs/{file.filename}"
        )
        
        # Process and index the file
        pages_indexed = _process_file(tmp_path, current_user["email"], file.filename, "cv")
        
        # Add level information to metadata
        for item in metadata:
            if (item["filename"] == file.filename and 
                item.get("doc_type") == "cv" and
                item.get("owner") == current_user["email"]):
                item["level"] = level
        
        _persist_stores()
        
        cv_id = str(uuid.uuid4())
        
        return UploadResponse(
            cvId=cv_id,
            fileName=file.filename,
            uploadedAt=datetime.utcnow().isoformat()
        )
    
    finally:
        # Clean up temporary file
        os.unlink(tmp_path)

# ── Skill Analysis Functions ──────────────────────────────────────────────
def _call_llm(prompt: str) -> str:
    """Call AWS Bedrock Claude for LLM analysis"""
    try:
        payload = {
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 4096,
            "messages": [
                {
                    "role": "user",
                    "content": [{"type": "text", "text": prompt}]
                }
            ]
        }
        
        response = bedrock.invoke_model(
            modelId=MODEL_ID,
            contentType="application/json",
            accept="application/json",
            body=json.dumps(payload)
        )
        
        body = json.loads(response["body"].read())
        return body["content"][0]["text"]
    
    except Exception as e:
        print(f"LLM call error: {e}")
        return f"Error in LLM analysis: {str(e)}"

def _extract_candidate_name(cv_text: str) -> str:
    """Extract candidate name from CV text using LLM"""
    prompt = f"""
Extract the candidate's full name from this CV/Resume text. 

CV Text:
{cv_text[:1000]}

Return only the candidate's name, nothing else. If no clear name is found, return "Name not found".
"""
    
    try:
        response = _call_llm(prompt)
        name = response.strip()
        # Clean up the response
        if name and len(name) < 100 and not name.lower().startswith("name not"):
            return name
        else:
            return "Name not found"
    except:
        return "Name not found"

def _extract_skills_from_text(text: str, doc_type: str) -> List[str]:
    """Extract skills from document text using LLM"""
    prompt = f"""
You are an expert HR analyst. Extract all technical skills, soft skills, and competencies from the following {doc_type}.

{doc_type.upper()} Content:
{text}

Instructions:
1. Extract ALL skills mentioned (technical, soft skills, tools, technologies, frameworks, etc.)
2. Return skills as a simple comma-separated list
3. Be comprehensive but avoid duplicates
4. Include both explicit skills and implied competencies
5. For experience levels, include them as skills (e.g., "5+ years Python")

Return only the comma-separated list of skills, nothing else.
"""
    
    response = _call_llm(prompt)
    # Parse the response to extract skills
    skills = [skill.strip() for skill in response.split(",") if skill.strip()]
    return skills

def _analyze_skill_match(jd_skills: List[str], cv_skills: List[str], jd_text: str, cv_text: str) -> Dict:
    """Perform detailed skill analysis between JD and CV"""
    prompt = f"""
You are an expert HR analyst performing skill gap analysis between a Job Description and a CV.

JOB DESCRIPTION SKILLS:
{', '.join(jd_skills)}

CV SKILLS:
{', '.join(cv_skills)}

JOB DESCRIPTION TEXT:
{jd_text[:2000]}...

CV TEXT:
{cv_text[:2000]}...

Perform a comprehensive analysis and provide:

1. MATCH SCORE (0-100): Overall compatibility score
2. SKILLS FOUND: Skills from JD that are present in CV
3. MISSING SKILLS: Critical skills from JD that are missing in CV
4. ADDITIONAL SKILLS: Valuable skills in CV not mentioned in JD
5. EXPERIENCE LEVEL MATCH: How well the experience level matches
6. DETAILED REASONING: Explain the scoring methodology and key factors

Format your response as:
MATCH_SCORE: [number]
SKILLS_FOUND: [comma-separated list]
MISSING_SKILLS: [comma-separated list]
ADDITIONAL_SKILLS: [comma-separated list]
EXPERIENCE_MATCH: [assessment]
REASONING: [detailed explanation]
"""
    
    return _call_llm(prompt)

def _generate_skill_report(jd_id: str, cv_analyses: List[Dict]) -> str:
    """Generate comprehensive skill analysis report"""
    report_sections = []
    
    report_sections.append("# SKILL ANALYSIS REPORT")
    report_sections.append(f"Generated on: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}")
    report_sections.append(f"Job Description ID: {jd_id}")
    report_sections.append(f"Number of CVs Analyzed: {len(cv_analyses)}")
    report_sections.append("")
    
    # Summary statistics
    if cv_analyses:
        scores = [analysis.get('match_score', 0) for analysis in cv_analyses]
        avg_score = sum(scores) / len(scores)
        max_score = max(scores)
        min_score = min(scores)
        
        report_sections.append("## SUMMARY STATISTICS")
        report_sections.append(f"Average Match Score: {avg_score:.1f}%")
        report_sections.append(f"Highest Match Score: {max_score}%")
        report_sections.append(f"Lowest Match Score: {min_score}%")
        report_sections.append("")
    
    # Individual CV analyses
    report_sections.append("## INDIVIDUAL CV ANALYSES")
    for i, analysis in enumerate(cv_analyses, 1):
        candidate_name = analysis.get('candidate_name', 'Name not found')
        cv_id = analysis.get('cv_id', 'Unknown')
        
        report_sections.append(f"### CV {i}: {candidate_name}")
        report_sections.append(f"CV ID: {cv_id}")
        report_sections.append(f"Match Score: {analysis.get('match_score', 0)}%")
        report_sections.append(f"Skills Found: {analysis.get('skills_found', 'None')}")
        report_sections.append(f"Missing Skills: {analysis.get('missing_skills', 'None')}")
        report_sections.append(f"Additional Skills: {analysis.get('additional_skills', 'None')}")
        report_sections.append(f"Experience Match: {analysis.get('experience_match', 'Unknown')}")
        report_sections.append("")
        report_sections.append("**Detailed Analysis:**")
        report_sections.append(analysis.get('reasoning', 'No detailed analysis available'))
        report_sections.append("")
        report_sections.append("---")
        report_sections.append("")
    
    # Recommendations
    if cv_analyses:
        report_sections.append("## RECOMMENDATIONS")
        top_candidates = sorted(cv_analyses, key=lambda x: x.get('match_score', 0), reverse=True)[:3]
        
        report_sections.append("### Top Candidates:")
        for i, candidate in enumerate(top_candidates, 1):
            candidate_name = candidate.get('candidate_name', 'Name not found')
            cv_id = candidate.get('cv_id', 'Unknown')
            match_score = candidate.get('match_score', 0)
            report_sections.append(f"{i}. {candidate_name} (CV {cv_id}) - {match_score}% match")
        
        report_sections.append("")
        report_sections.append("### Common Missing Skills:")
        all_missing = []
        for analysis in cv_analyses:
            missing = analysis.get('missing_skills', '').split(', ')
            all_missing.extend([skill.strip() for skill in missing if skill.strip()])
        
        from collections import Counter
        common_missing = Counter(all_missing).most_common(5)
        for skill, count in common_missing:
            report_sections.append(f"- {skill} (missing in {count}/{len(cv_analyses)} CVs)")
    
    return "\n".join(report_sections)

# ── Skill Analysis Endpoints ───────────────────────────────────────────────
@app.post("/v1/analysis")
async def perform_skill_analysis(
    request: AnalysisRequest,
    current_user: dict = Depends(get_current_user)
):
    """Perform skill analysis between JD and CVs"""
    analysis_id = str(uuid.uuid4())
    
    # Find JD content
    jd_content = []
    jd_skills = []
    
    for item in metadata:
        if (item.get("doc_type") == "jd" and 
            (item.get("doc_id") == request.jdId or 
             item["filename"] == request.jdId or
             item["filename"].replace(".pdf", "").replace(".docx", "") == request.jdId)):
            jd_content.append(item["text"])
    
    if not jd_content:
        raise HTTPException(404, "Job Description not found")
    
    jd_text = "\n\n".join(jd_content)
    
    # Extract skills from JD
    jd_skills = _extract_skills_from_text(jd_text, "job description")
    
    # Analyze each CV
    cv_analyses = []
    
    for cv_id in request.cvIds:
        # Find CV content
        cv_content = []
        cv_info = None
        
        for item in metadata:
            if (item.get("doc_type") == "cv" and 
                (item.get("doc_id") == cv_id or 
                 item["filename"] == cv_id or
                 item["filename"].replace(".pdf", "").replace(".docx", "") == cv_id)):
                cv_content.append(item["text"])
                if not cv_info:
                    cv_info = item
        
        if not cv_content:
            continue  # Skip if CV not found
        
        cv_text = "\n\n".join(cv_content)
        
        # Extract candidate name
        candidate_name = _extract_candidate_name(cv_text)
        
        # Extract skills from CV
        cv_skills = _extract_skills_from_text(cv_text, "resume/CV")
        
        # Perform skill matching analysis
        analysis_result = _analyze_skill_match(jd_skills, cv_skills, jd_text, cv_text)
        
        # Parse the LLM response
        analysis_lines = analysis_result.split('\n')
        parsed_analysis = {
            'cv_id': cv_id,
            'candidate_name': candidate_name,
            'match_score': 0,
            'skills_found': '',
            'missing_skills': '',
            'additional_skills': '',
            'experience_match': '',
            'reasoning': analysis_result
        }
        
        for line in analysis_lines:
            if line.startswith('MATCH_SCORE:'):
                try:
                    score_text = line.split(':', 1)[1].strip()
                    parsed_analysis['match_score'] = int(''.join(filter(str.isdigit, score_text)))
                except:
                    parsed_analysis['match_score'] = 0
            elif line.startswith('SKILLS_FOUND:'):
                parsed_analysis['skills_found'] = line.split(':', 1)[1].strip()
            elif line.startswith('MISSING_SKILLS:'):
                parsed_analysis['missing_skills'] = line.split(':', 1)[1].strip()
            elif line.startswith('ADDITIONAL_SKILLS:'):
                parsed_analysis['additional_skills'] = line.split(':', 1)[1].strip()
            elif line.startswith('EXPERIENCE_MATCH:'):
                parsed_analysis['experience_match'] = line.split(':', 1)[1].strip()
        
        cv_analyses.append(parsed_analysis)
    
    # Generate JSON-formatted analysis results
    jd_info = None
    for item in metadata:
        if (item.get("doc_type") == "jd" and 
            (item.get("doc_id") == request.jdId or 
             item["filename"] == request.jdId or
             item["filename"].replace(".pdf", "").replace(".docx", "") == request.jdId)):
            jd_info = item
            break
    
    jd_name = jd_info.get("title", "Unknown JD") if jd_info else "Unknown JD"
    
    # Create structured JSON response
    json_results = {
        "analysis_id": analysis_id,
        "total_candidates": len(cv_analyses),
        "jd_name": jd_name,
        "jd_id": request.jdId,
        "overall_analysis_notes": f"Skill analysis completed for {len(cv_analyses)} candidates against job description '{jd_name}'. Analysis includes skill matching, gap identification, and experience assessment using LLM-powered evaluation.",
        "candidates": []
    }
    
    # Add candidate details
    for analysis in cv_analyses:
        candidate_data = {
            "candidate_name": analysis.get('candidate_name', 'Name not found'),
            "cv_id": analysis.get('cv_id'),
            "match_score": analysis.get('match_score', 0),
            "skills_found": analysis.get('skills_found', ''),
            "missing_skills": analysis.get('missing_skills', ''),
            "additional_skills": analysis.get('additional_skills', ''),
            "experience_match": analysis.get('experience_match', ''),
            "detailed_reasoning": analysis.get('reasoning', '')
        }
        json_results["candidates"].append(candidate_data)
    
    # Sort candidates by match score (highest first)
    json_results["candidates"].sort(key=lambda x: x["match_score"], reverse=True)
    
    # Store analysis results (keeping both formats for backward compatibility)
    analysis_results[analysis_id] = {
        "analysisId": analysis_id,
        "timestamp": datetime.utcnow().isoformat(),
        "jdId": request.jdId,
        "cvIds": request.cvIds,
        "results": json_results,  # New JSON format
        "legacy_report": _generate_skill_report(request.jdId, cv_analyses),  # Keep old format for compatibility
        "cv_analyses": cv_analyses,
        "owner": current_user["email"]
    }
    
    save_analysis_results()
    
    return AnalysisResponse(
        analysisId=analysis_id,
        timestamp=datetime.utcnow().isoformat(),
        results=json_results  # Return the new JSON format
    )

@app.get("/v1/analysis/{analysisId}")
async def get_analysis_results(
    analysisId: str,
    current_user: dict = Depends(get_current_user)
):
    """Retrieve analysis results"""
    if analysisId not in analysis_results:
        raise HTTPException(404, "Analysis not found")
    
    analysis = analysis_results[analysisId]
    
    # Check access permissions
    if (analysis.get("owner") != current_user.get("email") and 
        current_user.get("role") != "admin"):
        raise HTTPException(403, "Access denied")
    
    return AnalysisResponse(
        analysisId=analysis["analysisId"],
        timestamp=analysis["timestamp"],
        results=analysis["results"]
    )

# ── Chat with Documents Endpoints ──────────────────────────────────────────
def _get_document_content(doc_id: str, doc_type: str) -> str:
    """Get document content by ID and type"""
    for item in metadata:
        if (item.get("doc_type") == doc_type and 
            (item.get("doc_id") == doc_id or 
             item["filename"] == doc_id or
             item["filename"].replace(".pdf", "").replace(".docx", "") == doc_id)):
            return item.get("text", "")
    return ""

def _generate_chat_response(jd_content: str, cv_content: str, conversation_history: List[Dict], user_message: str) -> str:
    """Generate chat response using LLM with document context"""
    
    # Prepare context with last 2 messages if they exist
    context_messages = ""
    if conversation_history:
        recent_messages = conversation_history[-2:]  # Last 2 messages
        for msg in recent_messages:
            context_messages += f"{msg['role'].title()}: {msg['content']}\n"
    
    prompt = f"""
You are an expert HR assistant helping users analyze job descriptions and CVs. You have access to the following documents:

JOB DESCRIPTION:
{jd_content[:2000]}...

CV/RESUME:
{cv_content[:2000]}...

CONVERSATION HISTORY (Last 2 messages):
{context_messages}

USER QUESTION:
{user_message}

Please provide a helpful, detailed response based on the documents and conversation context. You can:
- Answer questions about the candidate's qualifications
- Compare skills between JD and CV
- Provide insights about fit and gaps
- Suggest interview questions
- Explain specific aspects of either document
- Make recommendations

Keep your response conversational, informative, and focused on the documents provided.
"""
    
    try:
        response = bedrock.invoke_model(
            modelId=MODEL_ID,
            body=json.dumps({
                "anthropic_version": "bedrock-2023-05-31",
                "max_tokens": 1500,
                "messages": [{"role": "user", "content": prompt}]
            })
        )
        
        response_body = json.loads(response['body'].read())
        return response_body['content'][0]['text']
        
    except Exception as e:
        print(f"Error generating chat response: {e}")
        return "I apologize, but I'm having trouble processing your request right now. Please try again."

@app.post("/v1/chat")
async def chat_with_documents(
    request: ChatRequest,
    current_user: dict = Depends(get_current_user)
):
    """Start or continue a chat conversation about specific documents"""
    
    # Get document contents
    jd_content = _get_document_content(request.jdId, "jd")
    cv_content = _get_document_content(request.cvId, "cv")
    
    if not jd_content:
        raise HTTPException(404, "Job Description not found")
    if not cv_content:
        raise HTTPException(404, "CV not found")
    
    # Get or create conversation
    conversation_id = request.conversationId or str(uuid.uuid4())
    timestamp = datetime.utcnow().isoformat()
    
    # Load existing conversation or create new one
    if conversation_id in chat_conversations:
        conversation = chat_conversations[conversation_id]
        # Verify user has access to this conversation
        if conversation["userId"] != current_user["email"] and current_user["role"] != "admin":
            raise HTTPException(403, "Access denied to this conversation")
    else:
        # Create new conversation
        conversation = {
            "conversationId": conversation_id,
            "jdId": request.jdId,
            "cvId": request.cvId,
            "userId": current_user["email"],
            "messages": [],
            "createdAt": timestamp,
            "updatedAt": timestamp
        }
        chat_conversations[conversation_id] = conversation
    
    # Add user message to conversation
    user_msg = {
        "role": "user",
        "content": request.message,
        "timestamp": timestamp
    }
    conversation["messages"].append(user_msg)
    
    # Generate AI response
    ai_response = _generate_chat_response(
        jd_content, 
        cv_content, 
        conversation["messages"][:-1],  # Exclude the current message
        request.message
    )
    
    # Add AI response to conversation
    ai_msg = {
        "role": "assistant",
        "content": ai_response,
        "timestamp": datetime.utcnow().isoformat()
    }
    conversation["messages"].append(ai_msg)
    conversation["updatedAt"] = ai_msg["timestamp"]
    
    # Save conversations
    save_chat_conversations()
    
    return ChatResponse(
        conversationId=conversation_id,
        response=ai_response,
        timestamp=ai_msg["timestamp"]
    )

@app.get("/v1/chat/{conversationId}")
async def get_chat_conversation(
    conversationId: str,
    current_user: dict = Depends(get_current_user)
):
    """Retrieve a chat conversation"""
    if conversationId not in chat_conversations:
        raise HTTPException(404, "Conversation not found")
    
    conversation = chat_conversations[conversationId]
    
    # Check access permissions
    if (conversation["userId"] != current_user["email"] and 
        current_user["role"] != "admin"):
        raise HTTPException(403, "Access denied")
    
    return conversation

@app.get("/v1/chat")
async def list_chat_conversations(
    current_user: dict = Depends(get_current_user)
):
    """List all chat conversations for the current user"""
    user_conversations = []
    
    for conv_id, conversation in chat_conversations.items():
        if (conversation["userId"] == current_user["email"] or 
            current_user["role"] == "admin"):
            # Return summary without full message history
            summary = {
                "conversationId": conv_id,
                "jdId": conversation["jdId"],
                "cvId": conversation["cvId"],
                "messageCount": len(conversation["messages"]),
                "createdAt": conversation["createdAt"],
                "updatedAt": conversation["updatedAt"],
                "lastMessage": conversation["messages"][-1]["content"][:100] + "..." if conversation["messages"] else ""
            }
            user_conversations.append(summary)
    
    # Sort by most recent first
    user_conversations.sort(key=lambda x: x["updatedAt"], reverse=True)
    
    return {"conversations": user_conversations}

@app.delete("/v1/chat/{conversationId}")
async def delete_chat_conversation(
    conversationId: str,
    current_user: dict = Depends(get_current_user)
):
    """Delete a chat conversation"""
    if conversationId not in chat_conversations:
        raise HTTPException(404, "Conversation not found")
    
    conversation = chat_conversations[conversationId]
    
    # Check permissions
    if (conversation["userId"] != current_user["email"] and 
        current_user["role"] != "admin"):
        raise HTTPException(403, "Access denied")
    
    del chat_conversations[conversationId]
    save_chat_conversations()
    
    return {"message": "Conversation deleted successfully"}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

