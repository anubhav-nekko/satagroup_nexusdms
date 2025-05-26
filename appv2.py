# appv2.py

import os, json, tempfile, pickle, fitz, faiss, boto3
from pathlib import Path
from typing import List, Dict, Tuple, Optional
from datetime import datetime
import numpy as np
import pandas as pd
from fastapi import FastAPI, UploadFile, File, HTTPException, Body, Query
from pydantic import BaseModel
from sentence_transformers import SentenceTransformer
from docx import Document
from pptx import Presentation

# ── CONFIG ──────────────────────────────────────────────────────────────────
REGION   = "us-east-1"
BUCKET   = "satagroup-test"
MODEL_ID = "arn:aws:bedrock:us-east-1:343218220592:inference-profile/us.anthropic.claude-3-7-sonnet-20250219-v1:0"

INDEX_F  = "faiss_index.bin"
META_F   = "metadata_store.pkl"
EMB_DIM  = 768

# ── AWS clients ─────────────────────────────────────────────────────────────
s3       = boto3.client("s3", region_name=REGION)
textract = boto3.client("textract", region_name=REGION)
bedrock  = boto3.client("bedrock-runtime", region_name=REGION)

# ── Globals ─────────────────────────────────────────────────────────────────
app         = FastAPI(title="Document-RAG Gateway", version="1.2.0")
embed_model = SentenceTransformer("sentence-transformers/all-mpnet-base-v2")

# These will be loaded or created at startup:
faiss_index : faiss.Index
metadata    : List[Dict]

# ── Utility - load / save stores ────────────────────────────────────────────
def _load_stores() -> None:
    """Load FAISS index and metadata from S3 if they exist."""
    global faiss_index, metadata
    try:
        s3.download_file(BUCKET, INDEX_F, INDEX_F)
        s3.download_file(BUCKET, META_F,  META_F)
        faiss_index = faiss.read_index(INDEX_F)
        metadata    = pickle.load(open(META_F, "rb"))
    except Exception:
        # If not found or any error, start with empty.
        faiss_index = faiss.IndexFlatL2(EMB_DIM)
        metadata    = []

def _persist_stores() -> None:
    """Write FAISS index + metadata to local disk and then upload to S3."""
    faiss.write_index(faiss_index, INDEX_F)
    pickle.dump(metadata, open(META_F, "wb"))
    s3.upload_file(INDEX_F, BUCKET, INDEX_F)
    s3.upload_file(META_F, BUCKET, META_F)

def _embed(text: str) -> np.ndarray:
    """Use your sentence-transformers model to get normalized embeddings."""
    return embed_model.encode(text, normalize_embeddings=True)

def _ocr_png(png_bytes: bytes) -> str:
    """Use Textract to OCR PNG images (including PDF pages rendered as PNG)."""
    resp = textract.detect_document_text(Document={"Bytes": png_bytes})
    lines = [b["Text"] for b in resp.get("Blocks", []) if b["BlockType"] == "LINE"]
    return "\n".join(lines)

# ── Core extractor for a single file ────────────────────────────────────────
def _process_file(path: Path, owner: str, real_name: str) -> int:
    """
    Read the given file, chunk/ocr as needed, embed, store into FAISS index
    and keep track in `metadata`. Returns the number of non-empty chunks.
    """
    ext = path.suffix.lower()
    chunks: List[Tuple[str,int]] = []

    if ext == ".pdf":
        doc = fitz.open(path)
        for pg in doc:
            png = pg.get_pixmap(dpi=300).tobytes("png")
            text = _ocr_png(png)
            if text.strip():
                chunks.append((text, pg.number + 1))

    elif ext in {".jpg", ".jpeg", ".png"}:
        text = _ocr_png(path.read_bytes())
        if text.strip():
            chunks.append((text, 1))

    elif ext in {".doc", ".docx"}:
        full = "\n".join(p.text for p in Document(path).paragraphs if p.text)
        # chunk ~1000 characters each
        for i in range(0, len(full), 1000):
            chunk = full[i:i+1000]
            if chunk.strip():
                chunks.append((chunk, i//1000 + 1))

    elif ext == ".pptx":
        prs = Presentation(path)
        for idx, slide in enumerate(prs.slides, start=1):
            slide_text = "\n".join(s.text for s in slide.shapes if hasattr(s, "text"))
            if slide_text.strip():
                chunks.append((slide_text, idx))

    elif ext in {".csv", ".xlsx"}:
        if ext == ".csv":
            df = pd.read_csv(path)
        else:
            df = pd.read_excel(path)
        # chunk 50 rows at a time
        for i in range(0, len(df), 50):
            chunk = df.iloc[i:i+50].to_string(index=False)
            if chunk.strip():
                chunks.append((chunk, i//50 + 1))

    else:
        raise HTTPException(400, f"Unsupported file type {ext}")

    # Embed & store
    for txt, page_num in chunks:
        vec = _embed(txt).reshape(1, -1)
        faiss_index.add(vec)
        metadata.append({
            "filename": real_name,
            "page":     page_num,
            "text":     txt,
            "owner":    owner,
            "uploaded": datetime.utcnow().isoformat()
        })
    return len(chunks)

# ── Models / Schemas ───────────────────────────────────────────────────────
class UploadAck(BaseModel):
    status: str
    filename: str
    pages_indexed: int

class BulkUploadAck(BaseModel):
    status: str
    files: List[UploadAck]

class QueryBody(BaseModel):
    selected_files: List[str]
    selected_page_ranges: Dict[str, Tuple[int, int]] = {}
    prompt: str
    top_k: int = 20
    last_messages: List[str] = []
    web_search: bool = False

class Answer(BaseModel):
    answer: str
    context: List[Dict]

# ── Startup ────────────────────────────────────────────────────────────────
@app.on_event("startup")
def _startup():
    _load_stores()

# ── Original single-file Upload ────────────────────────────────────────────
@app.post("/upload_document", response_model=UploadAck)
async def upload_document(file: UploadFile = File(...), owner: str = ""):
    if not owner:
        raise HTTPException(400, "Missing owner field.")

    fn = file.filename or ""
    if not fn:
        raise HTTPException(400, "No filename provided.")

    # Save file to a temp location
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=Path(fn).suffix)
    tmp.write(await file.read())
    tmp.close()

    # Store original in S3
    with open(tmp.name, "rb") as fh:
        s3.upload_fileobj(fh, BUCKET, fn)

    # Process for embeddings
    try:
        pages_indexed = _process_file(Path(tmp.name), owner, fn)
        _persist_stores()
    finally:
        os.remove(tmp.name)

    return UploadAck(status="done", filename=fn, pages_indexed=pages_indexed)

# ── 1) Bulk Upload ─────────────────────────────────────────────────────────
@app.post("/bulk_upload", response_model=BulkUploadAck)
async def bulk_upload(
    files: List[UploadFile] = File(...),
    owner: str = Query(..., description="Owner or user ID uploading the files"),
):
    """
    Accepts multiple files in one request.  
    Returns a list with the status and pages indexed for each file.
    """
    if not owner:
        raise HTTPException(400, "Missing owner field.")

    results = []
    for file in files:
        fn = file.filename or ""
        if not fn:
            # skip
            continue
        
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=Path(fn).suffix)
        tmp.write(await file.read())
        tmp.close()

        # Upload the original file to S3
        with open(tmp.name, "rb") as fh:
            s3.upload_fileobj(fh, BUCKET, fn)

        # Process + embed
        try:
            pages = _process_file(Path(tmp.name), owner, fn)
            results.append(UploadAck(status="done", filename=fn, pages_indexed=pages))
        except Exception as e:
            results.append(UploadAck(status=f"failed: {str(e)}", filename=fn, pages_indexed=0))
        finally:
            os.remove(tmp.name)

    # Persist once after all files are processed
    _persist_stores()

    return BulkUploadAck(status="bulk_done", files=results)

# ── 2) Conversation Log Retrieval ──────────────────────────────────────────
def _load_chat_history() -> dict:
    """
    Example helper that loads a JSON file from S3 named `chat_history.json`.
    This is the same logic you'd have used in Streamlit. Adjust as needed.
    """
    blob_name = "chat_history.json"
    try:
        # Check if object exists
        s3.head_object(Bucket=BUCKET, Key=blob_name)
    except:
        # No chat history in S3
        return {}

    # If found, download and load
    tmp_file = tempfile.NamedTemporaryFile(delete=False)
    s3.download_file(BUCKET, blob_name, tmp_file.name)
    tmp_file.close()

    data = {}
    with open(tmp_file.name, "r", encoding="utf-8") as f:
        data = json.load(f)
    os.remove(tmp_file.name)
    return data if isinstance(data, dict) else {}

@app.get("/conversation_log/{username}")
def get_conversation_log(username: str) -> dict:
    """
    Return the conversation history for a particular user.
    Structure depends on how you store it in `_load_chat_history()`.
    """
    history = _load_chat_history()
    user_logs = history.get(username, [])
    return {"username": username, "history": user_logs}

# ── 3) Query Documents (JSON-based) ────────────────────────────────────────
@app.post("/query_documents", response_model=Answer)
def query_docs(body: QueryBody):
    """
    Similar to your /query_documents_with_page_range, but named more generally.
    Returns JSON with the 'answer' and the 'context' hits used.
    """
    if faiss_index.ntotal == 0:
        raise HTTPException(404, "Index is empty; no documents to search.")

    q_vec = _embed(body.prompt).reshape(1, -1)

    # We'll over-fetch, then filter by file + page ranges
    topN = min(body.top_k * 5, faiss_index.ntotal)
    D, I = faiss_index.search(q_vec, topN)

    # Filter & pick best hits
    hits = []
    for dist, idx in zip(D[0], I[0]):
        meta = metadata[idx]
        f    = meta["filename"]
        p    = meta["page"]

        # If user selected a subset of files, skip if not in that subset
        if body.selected_files and f not in body.selected_files:
            continue

        # If page ranges are specified, skip if out of range
        lo, hi = body.selected_page_ranges.get(f, (1, 10**9))
        if not(lo <= p <= hi):
            continue

        hits.append({"score": float(dist), **meta})
        if len(hits) >= body.top_k:
            break

    # Build the final prompt for the LLM
    # (You can incorporate last_messages, web_search, etc.)
    context_text = json.dumps({"hits": hits, "chat": body.last_messages}, indent=2)
    llm_prompt   = body.prompt + "\n\n# Context\n" + context_text

    # Call Anthropic / Claude via Bedrock
    payload = {
        "anthropic_version": "bedrock-2023-05-31",
        "max_tokens": 4096,
        "messages": [
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": llm_prompt
                    }
                ]
            }
        ],
    }
    out    = bedrock.invoke_model(
                modelId=MODEL_ID,
                contentType="application/json",
                accept="application/json",
                body=json.dumps(payload),
             )
    answer = json.loads(out["body"].read())["content"][0]["text"]

    return Answer(answer=answer, context=hits)

# ── 4) List Available Files ────────────────────────────────────────────────
@app.get("/list_files", response_model=List[str])
def list_files(owner: Optional[str] = None) -> List[str]:
    """
    Return a distinct list of filenames from `metadata`.
    Optionally filter by `owner` if provided.
    """
    if not metadata:
        return []

    unique_files = set()
    for rec in metadata:
        if owner:
            if rec.get("owner") == owner:
                unique_files.add(rec["filename"])
        else:
            unique_files.add(rec["filename"])

    return sorted(unique_files)
